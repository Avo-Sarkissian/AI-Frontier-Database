"""
Build the AI Frontier presentation deck.

Theme matches the live dashboard: near-black canvas (#0a0a0a), elevated cards
(#111111), cyan accent (#00d4ff), Inter typography, generous whitespace,
high data-density layouts. 16:9, ~10 slides, paced for a 10-minute talk
with a substantial live demo at the end.

Run:  python3 build_presentation.py
Out:  AI_Frontier_Presentation.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ── Theme tokens (mirrors assets/style.css) ─────────────────────────────────
BG          = RGBColor(0x0A, 0x0A, 0x0A)
BG_CARD     = RGBColor(0x11, 0x11, 0x11)
BG_ELEVATED = RGBColor(0x16, 0x16, 0x16)
BORDER      = RGBColor(0x26, 0x26, 0x26)
TEXT_1      = RGBColor(0xED, 0xED, 0xED)
TEXT_2      = RGBColor(0x88, 0x88, 0x88)
TEXT_3      = RGBColor(0x66, 0x66, 0x66)
TEXT_4      = RGBColor(0x44, 0x44, 0x44)
ACCENT      = RGBColor(0x00, 0xD4, 0xFF)
ACCENT_DIM  = RGBColor(0x0E, 0x35, 0x44)
GREEN       = RGBColor(0x22, 0xC5, 0x5E)
PURPLE      = RGBColor(0xC0, 0x84, 0xFC)
ORANGE      = RGBColor(0xFB, 0x92, 0x3C)

FONT = "Inter"
FONT_FALLBACK = "Helvetica Neue"

# Slide geometry: 16:9 widescreen, 13.333" × 7.5"
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


# ── Helpers ─────────────────────────────────────────────────────────────────
def set_bg(slide, color=BG):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    bg.shadow.inherit = False
    return bg


def add_rect(slide, x, y, w, h, fill=BG_CARD, line=None, line_w=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = line_w or Pt(0.75)
    shp.shadow.inherit = False
    return shp


def add_line(slide, x1, y1, x2, y2, color=BORDER, width=Pt(0.75)):
    ln = slide.shapes.add_connector(1, x1, y1, x2, y2)
    ln.line.color.rgb = color
    ln.line.width = width
    return ln


def add_text(slide, x, y, w, h, text, *,
             size=12, color=TEXT_1, bold=False, font=FONT,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             tracking=None):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    p.space_before = Pt(0)
    p.space_after = Pt(0)
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    if tracking is not None:
        # tracking in 1/100 of a point
        rPr = run._r.get_or_add_rPr()
        rPr.set("spc", str(tracking))
    return tb


def add_paragraphs(slide, x, y, w, h, lines, *,
                   size=12, color=TEXT_2, bold=False,
                   font=FONT, leading=Pt(6), align=PP_ALIGN.LEFT):
    """Add a multi-paragraph block. Each `line` is a string OR a dict
    {text, size, color, bold}."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_before = Pt(0)
        p.space_after = leading
        if isinstance(ln, dict):
            run = p.add_run()
            run.text = ln.get("text", "")
            run.font.name = ln.get("font", font)
            run.font.size = Pt(ln.get("size", size))
            run.font.bold = ln.get("bold", bold)
            run.font.color.rgb = ln.get("color", color)
        else:
            run = p.add_run()
            run.text = ln
            run.font.name = font
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.color.rgb = color
    return tb


def add_chip(slide, x, y, label, *, color=ACCENT, fill=ACCENT_DIM, size=9):
    """Pill-shaped tag matching the dashboard's header-badge style."""
    h = Inches(0.28)
    # Approximate width from label length
    w = Inches(0.18 + 0.085 * len(label))
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shp.adjustments[0] = 0.5
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = color
    shp.line.width = Pt(0.75)
    tf = shp.text_frame
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = label.upper()
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = True
    run.font.color.rgb = color
    rPr = run._r.get_or_add_rPr()
    rPr.set("spc", "120")
    return shp, w


def add_header(slide, eyebrow, title, page_num):
    """Sticky-style header — small uppercase eyebrow over a tight title."""
    # Eyebrow chip
    add_chip(slide, Inches(0.6), Inches(0.45), eyebrow)

    # Page number, top right
    add_text(slide, SLIDE_W - Inches(1.2), Inches(0.5), Inches(0.6), Inches(0.3),
             f"{page_num:02d} / 10", size=9, color=TEXT_3,
             align=PP_ALIGN.RIGHT, tracking=120, bold=True)

    # Title row
    add_text(slide, Inches(0.6), Inches(0.85), Inches(12), Inches(0.7),
             title, size=30, color=TEXT_1, bold=True)

    # Divider line under header
    y = Inches(1.6)
    add_line(slide, Inches(0.6), y, SLIDE_W - Inches(0.6), y, color=BORDER, width=Pt(0.75))


def add_footer(slide):
    add_text(slide, Inches(0.6), SLIDE_H - Inches(0.4), Inches(6),
             Inches(0.3), "AI FRONTIER  ·  EECE 5642  ·  SPRING 2026",
             size=8, color=TEXT_3, tracking=160, bold=True)
    add_text(slide, SLIDE_W - Inches(4), SLIDE_H - Inches(0.4),
             Inches(3.4), Inches(0.3),
             "github.com/Avo-Sarkissian/AI-Frontier-Database",
             size=8, color=TEXT_3, tracking=80, align=PP_ALIGN.RIGHT)


def stat_card(slide, x, y, w, h, value, label, *, accent=False):
    add_rect(slide, x, y, w, h, fill=BG_CARD)
    add_text(slide, x + Inches(0.25), y + Inches(0.30), w - Inches(0.5),
             Inches(0.7), value, size=28, bold=True,
             color=ACCENT if accent else TEXT_1)
    add_text(slide, x + Inches(0.25), y + Inches(0.95), w - Inches(0.5),
             Inches(0.3), label, size=9, color=TEXT_3,
             tracking=120, bold=True)


# ════════════════════════════════════════════════════════════════════════════
#  BUILD
# ════════════════════════════════════════════════════════════════════════════
prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]


# ── Slide 1 — TITLE ────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
set_bg(s)

# Subtle accent bar across the top
add_rect(s, 0, 0, SLIDE_W, Inches(0.06), fill=ACCENT)

# Cyan dot + brand mark
add_rect(s, Inches(0.6), Inches(0.55), Inches(0.12), Inches(0.12), fill=ACCENT)
add_text(s, Inches(0.85), Inches(0.5), Inches(6), Inches(0.3),
         "AI FRONTIER", size=11, bold=True, color=TEXT_1, tracking=200)

add_chip(s, SLIDE_W - Inches(2.0), Inches(0.5), "LIVE DEMO  ●", color=ACCENT, fill=ACCENT_DIM)

# Title block, vertically centered
add_text(s, Inches(0.6), Inches(2.2), Inches(12), Inches(0.4),
         "DATA VISUALIZATION  ·  EECE 5642", size=10, color=TEXT_3,
         bold=True, tracking=200)

add_text(s, Inches(0.6), Inches(2.7), Inches(12.2), Inches(1.4),
         "Mapping the AI Model Landscape", size=52, bold=True,
         color=TEXT_1)

add_text(s, Inches(0.6), Inches(4.05), Inches(12.2), Inches(0.7),
         "An interactive dashboard comparing 255+ language models on",
         size=20, color=TEXT_2)
add_text(s, Inches(0.6), Inches(4.45), Inches(12.2), Inches(0.7),
         "cost, speed, and intelligence — updated hourly.",
         size=20, color=TEXT_2)

# Bottom stat bar — mirrors the dashboard's stat strip
y0 = Inches(5.7)
add_line(s, Inches(0.6), y0, SLIDE_W - Inches(0.6), y0, color=BORDER)
stats = [
    ("255+", "MODELS TRACKED", False),
    ("29",   "PROVIDERS",      False),
    ("11",   "INTERACTIVE TABS", False),
    ("HOURLY", "DATA REFRESH",  True),
]
sw = (SLIDE_W - Inches(1.2)) / len(stats)
for i, (v, l, accent) in enumerate(stats):
    add_text(s, Inches(0.6) + sw * i, y0 + Inches(0.2), sw, Inches(0.6),
             v, size=26, bold=True, color=ACCENT if accent else TEXT_1)
    add_text(s, Inches(0.6) + sw * i, y0 + Inches(0.85), sw, Inches(0.3),
             l, size=9, color=TEXT_3, tracking=120, bold=True)

add_text(s, Inches(0.6), SLIDE_H - Inches(0.45), Inches(8), Inches(0.3),
         "Avo Sarkissian  ·  Northeastern University  ·  Spring 2026",
         size=10, color=TEXT_3)


# ── Slide 2 — WHY THIS MATTERS ─────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
set_bg(s)
add_header(s, "01 · CONTEXT", "Why these tradeoffs matter", 2)
add_footer(s)

add_text(s, Inches(0.6), Inches(1.85), Inches(12), Inches(0.5),
         "Dozens of new models ship every week. Picking the wrong one is the difference between a $5 and a $5,000 monthly bill.",
         size=14, color=TEXT_2)

# Four audience cards
cards = [
    ("DEVELOPERS",
     "Match each task to the cheapest model that still meets quality. A coding agent doesn't need GPT-5-class reasoning to rename a variable.",
     ACCENT),
    ("RESEARCHERS",
     "Track frontier movement over time. The Pareto frontier shifts every few weeks — what was state-of-the-art is now mid-tier.",
     PURPLE),
    ("BUSINESSES",
     "Project monthly token spend before signing a contract. Two models with identical benchmarks can differ 50× in price.",
     GREEN),
    ("STUDENTS  ·  HOBBYISTS",
     "Decide what runs on your laptop vs. what needs an API key. Open-weight models close the gap faster than most people think.",
     ORANGE),
]
gap = Inches(0.18)
cw = (SLIDE_W - Inches(1.2) - gap * 3) / 4
ch = Inches(3.3)
cy = Inches(2.7)
for i, (title, body, col) in enumerate(cards):
    cx = Inches(0.6) + (cw + gap) * i
    add_rect(s, cx, cy, cw, ch, fill=BG_CARD)
    # accent bar at top
    add_rect(s, cx, cy, cw, Inches(0.05), fill=col)
    add_text(s, cx + Inches(0.25), cy + Inches(0.3), cw - Inches(0.5),
             Inches(0.4), title, size=10, bold=True, color=col, tracking=160)
    add_text(s, cx + Inches(0.25), cy + Inches(0.85), cw - Inches(0.5),
             Inches(2.3), body, size=12, color=TEXT_1)

add_text(s, Inches(0.6), Inches(6.4), Inches(12), Inches(0.4),
         "AI Frontier collapses every comparison into one lens — so the answer is a glance, not a spreadsheet.",
         size=12, color=TEXT_3)


# ── Slide 3 — DATA SOURCE & FRESHNESS ──────────────────────────────────────
s = prs.slides.add_slide(BLANK)
set_bg(s)
add_header(s, "02 · DATA SOURCE", "One API, refreshed every hour", 3)
add_footer(s)

# Left: pipeline diagram
lx = Inches(0.6); ly = Inches(2.0); lw = Inches(7.2); lh = Inches(4.6)
add_rect(s, lx, ly, lw, lh, fill=BG_CARD)
add_text(s, lx + Inches(0.3), ly + Inches(0.25), lw - Inches(0.6), Inches(0.3),
         "INGESTION PIPELINE", size=9, bold=True, color=TEXT_3, tracking=160)

# Pipeline boxes — vertical flow
steps = [
    ("Artificial Analysis API",
     "host-models/performance endpoint  ·  ~20 MB JSON  ·  800+ host records",
     ACCENT),
    ("Background scraper thread",
     "data/scraper.py  ·  fires on startup, then every 3,600 s",
     TEXT_1),
    ("Aggregate by (model, lab)",
     "Keep cheapest price + fastest speed across all hosts",
     TEXT_1),
    ("CSV cache + daily snapshot",
     "data/raw/aa_models.csv  +  data/raw/history/aa_models_YYYY-MM-DD.csv",
     TEXT_1),
    ("Dash app reloads on stale mtime",
     "Charts re-render automatically — no page refresh needed",
     GREEN),
]
sy = ly + Inches(0.7)
sh = Inches(0.7)
for i, (title, sub, col) in enumerate(steps):
    yy = sy + (sh + Inches(0.06)) * i
    add_rect(s, lx + Inches(0.4), yy, lw - Inches(0.8), sh, fill=BG_ELEVATED)
    # left accent stripe
    add_rect(s, lx + Inches(0.4), yy, Inches(0.06), sh, fill=col)
    add_text(s, lx + Inches(0.6), yy + Inches(0.1),
             lw - Inches(1.2), Inches(0.3), title, size=12, bold=True, color=TEXT_1)
    add_text(s, lx + Inches(0.6), yy + Inches(0.35),
             lw - Inches(1.2), Inches(0.3), sub, size=10, color=TEXT_3)

# Right: freshness metrics
rx = Inches(8.1); ry = Inches(2.0); rw = Inches(4.65); rh = Inches(4.6)
add_rect(s, rx, ry, rw, rh, fill=BG_CARD)
add_text(s, rx + Inches(0.3), ry + Inches(0.25), rw - Inches(0.6), Inches(0.3),
         "FRESHNESS", size=9, bold=True, color=TEXT_3, tracking=160)

# Big number 1
add_text(s, rx + Inches(0.3), ry + Inches(0.7), rw - Inches(0.6), Inches(0.7),
         "1 hr", size=40, bold=True, color=ACCENT)
add_text(s, rx + Inches(0.3), ry + Inches(1.6), rw - Inches(0.6), Inches(0.3),
         "Background scrape interval", size=10, color=TEXT_2)

add_line(s, rx + Inches(0.3), ry + Inches(2.1),
         rx + rw - Inches(0.3), ry + Inches(2.1), color=BORDER)

# Big number 2
add_text(s, rx + Inches(0.3), ry + Inches(2.25), rw - Inches(0.6), Inches(0.7),
         "30+", size=40, bold=True, color=TEXT_1)
add_text(s, rx + Inches(0.3), ry + Inches(3.15), rw - Inches(0.6), Inches(0.3),
         "Daily snapshots since 2026-02-25", size=10, color=TEXT_2)

add_line(s, rx + Inches(0.3), ry + Inches(3.65),
         rx + rw - Inches(0.3), ry + Inches(3.65), color=BORDER)

# Big number 3
add_text(s, rx + Inches(0.3), ry + Inches(3.8), rw - Inches(0.6), Inches(0.7),
         "255+", size=40, bold=True, color=TEXT_1)
add_text(s, rx + Inches(0.3), ry + Inches(4.05), rw - Inches(0.6), Inches(0.3),
         " ", size=10, color=TEXT_2)
add_text(s, rx + Inches(0.3), ry + Inches(4.2), rw - Inches(0.6), Inches(0.3),
         "Models in latest snapshot", size=10, color=TEXT_2)


# ── Slide 4 — INTELLIGENCE INDEX ───────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
set_bg(s)
add_header(s, "03 · METHODOLOGY", "How the Intelligence Index is built", 4)
add_footer(s)

add_text(s, Inches(0.6), Inches(1.85), Inches(12), Inches(0.5),
         "A composite score from Artificial Analysis — averaged across seven evaluations covering reasoning, coding, and knowledge.",
         size=13, color=TEXT_2)

# Seven-eval grid
evals = [
    ("MMLU-Pro",      "Multitask knowledge",       "Reasoning"),
    ("GPQA Diamond",  "Graduate science",          "Reasoning"),
    ("Humanity's Last Exam", "Frontier reasoning",  "Reasoning"),
    ("LiveCodeBench", "Code generation",            "Coding"),
    ("SciCode",       "Scientific code",            "Coding"),
    ("AIME 2025",     "Competition math",           "Math"),
    ("MATH-500",      "Quantitative reasoning",     "Math"),
]
gap = Inches(0.12)
cols = 4
cw = (SLIDE_W - Inches(1.2) - gap * (cols - 1)) / cols
ch = Inches(1.15)
ey = Inches(2.55)
for i, (name, sub, cat) in enumerate(evals):
    r = i // cols
    c = i % cols
    cx = Inches(0.6) + (cw + gap) * c
    cy = ey + (ch + gap) * r
    add_rect(s, cx, cy, cw, ch, fill=BG_CARD)
    add_text(s, cx + Inches(0.2), cy + Inches(0.18),
             cw - Inches(0.4), Inches(0.3), cat.upper(),
             size=8, bold=True, color=ACCENT, tracking=140)
    add_text(s, cx + Inches(0.2), cy + Inches(0.45),
             cw - Inches(0.4), Inches(0.4), name,
             size=14, bold=True, color=TEXT_1)
    add_text(s, cx + Inches(0.2), cy + Inches(0.78),
             cw - Inches(0.4), Inches(0.3), sub,
             size=10, color=TEXT_3)

# Formula card spans the empty 8th cell
cx = Inches(0.6) + (cw + gap) * 3
cy = ey + (ch + gap) * 1
add_rect(s, cx, cy, cw, ch, fill=BG_ELEVATED, line=ACCENT, line_w=Pt(1))
add_text(s, cx + Inches(0.2), cy + Inches(0.18),
         cw - Inches(0.4), Inches(0.3), "FORMULA",
         size=8, bold=True, color=ACCENT, tracking=140)
add_text(s, cx + Inches(0.2), cy + Inches(0.42),
         cw - Inches(0.4), Inches(0.5),
         "mean(7 evals)", size=15, bold=True, color=TEXT_1)
add_text(s, cx + Inches(0.2), cy + Inches(0.78),
         cw - Inches(0.4), Inches(0.3),
         "Normalised 0–100", size=10, color=TEXT_3)

# Bottom note
ny = Inches(5.6)
add_rect(s, Inches(0.6), ny, SLIDE_W - Inches(1.2), Inches(1.3), fill=BG_CARD)
add_text(s, Inches(0.85), ny + Inches(0.2), Inches(12), Inches(0.3),
         "WHY A COMPOSITE?", size=9, bold=True, color=TEXT_3, tracking=140)
add_text(s, Inches(0.85), ny + Inches(0.5), Inches(12), Inches(0.4),
         "Single benchmarks are gameable. A model that crushes MMLU but flunks coding is misleading on its own.",
         size=12, color=TEXT_1)
add_text(s, Inches(0.85), ny + Inches(0.85), Inches(12), Inches(0.4),
         "Averaging seven independent evals reduces variance and rewards generalists over benchmark specialists.",
         size=12, color=TEXT_2)


# ── Slide 5 — TECH STACK ───────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
set_bg(s)
add_header(s, "04 · TOOLS", "Built entirely in Python", 5)
add_footer(s)

add_text(s, Inches(0.6), Inches(1.85), Inches(12), Inches(0.4),
         "No JS frameworks, no build step. The dashboard is a single Dash app served by Gunicorn.",
         size=13, color=TEXT_2)

libs = [
    ("Dash 4",     "Reactive web framework",
     "Tabs, callbacks, layout — turns the chart functions into an interactive site without writing HTML.", ACCENT),
    ("Plotly 6",   "Charting engine",
     "All 19 chart modules use plotly.graph_objects directly — fine control over traces, hover, theming.", PURPLE),
    ("Pandas",     "Data wrangling",
     "Loads the CSV cache, applies filters, computes Pareto frontier, value scores, percentile presets.", GREEN),
    ("NumPy",      "Numerical ops",
     "log-scaling, Pearson r between log10(price) and quality, bubble-size normalisation.", ORANGE),
    ("Requests",   "Live scraper",
     "Plain HTTP GET against the AA API — no Selenium, no Playwright, no browser overhead.", ACCENT),
    ("python-pptx", "This deck",
     "The slides you're looking at are generated programmatically — same color tokens as the website.", PURPLE),
]
gap = Inches(0.18)
cols = 3
cw = (SLIDE_W - Inches(1.2) - gap * (cols - 1)) / cols
ch = Inches(2.05)
ly = Inches(2.45)
for i, (name, sub, body, col) in enumerate(libs):
    r = i // cols
    c = i % cols
    cx = Inches(0.6) + (cw + gap) * c
    cy = ly + (ch + gap) * r
    add_rect(s, cx, cy, cw, ch, fill=BG_CARD)
    add_rect(s, cx, cy, Inches(0.06), ch, fill=col)
    add_text(s, cx + Inches(0.3), cy + Inches(0.25), cw - Inches(0.5),
             Inches(0.4), name, size=16, bold=True, color=TEXT_1)
    add_text(s, cx + Inches(0.3), cy + Inches(0.62), cw - Inches(0.5),
             Inches(0.3), sub.upper(), size=8, bold=True,
             color=col, tracking=140)
    add_text(s, cx + Inches(0.3), cy + Inches(0.95), cw - Inches(0.5),
             Inches(1.0), body, size=11, color=TEXT_2)


# ── Slide 6 — VISUALIZATION TYPES ──────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
set_bg(s)
add_header(s, "05 · CHART INVENTORY", "What's in the dashboard", 6)
add_footer(s)

add_text(s, Inches(0.6), Inches(1.85), Inches(12), Inches(0.4),
         "Eleven tabs, nineteen chart modules. Each chart is matched to its analytical task — never a chart for its own sake.",
         size=13, color=TEXT_2)

# Header row
hy = Inches(2.55)
add_rect(s, Inches(0.6), hy, SLIDE_W - Inches(1.2), Inches(0.4), fill=BG_ELEVATED)
cols_def = [
    ("CHART TYPE",     Inches(0.85),  Inches(2.5)),
    ("WHERE",          Inches(3.45),  Inches(2.4)),
    ("ENCODES",        Inches(5.95),  Inches(3.6)),
    ("WHY THIS CHART", Inches(9.65),  Inches(3.5)),
]
for label, x, w in cols_def:
    add_text(s, x, hy + Inches(0.10), w, Inches(0.3), label,
             size=8, bold=True, color=TEXT_3, tracking=140)

rows = [
    ("Bubble scatter",   "Overview",       "Price · Quality · Speed",     "Three quantitative axes at once"),
    ("Pareto frontier",  "Overview overlay", "Optimal cost/quality tier",  "Highlights non-dominated models"),
    ("Treemap",          "Landscape",      "Provider count + avg quality", "Compare hierarchical magnitudes"),
    ("Horizontal bars",  "Rankings",       "Score · Value · Speed",       "Best for ordered comparisons"),
    ("Radar",            "Compare",        "5 metrics, up to 5 models",   "Profile shape comparison"),
    ("Slope / bump chart", "Insights",     "Rank evolution over time",    "Shows movement, not just current rank"),
    ("Faceted scatter",  "Image Gen",      "ELO × Price by provider",     "Small multiples reduce overplotting"),
    ("Quadrant scatter", "Run Local",      "VRAM fit × Speed",            "Decision space split by hardware"),
]
ry0 = hy + Inches(0.5)
rh = Inches(0.45)
for i, row in enumerate(rows):
    yy = ry0 + rh * i
    if i % 2 == 1:
        add_rect(s, Inches(0.6), yy, SLIDE_W - Inches(1.2), rh, fill=BG_CARD)
    for j, (label, x, w) in enumerate(cols_def):
        text = row[j]
        col = TEXT_1 if j == 0 else (ACCENT if j == 1 else TEXT_2)
        bold = (j == 0)
        size = 11
        add_text(s, x, yy + Inches(0.13), w, Inches(0.3),
                 text, size=size, bold=bold, color=col)


# ── Slide 7 — DATA VIZ PRINCIPLES ──────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
set_bg(s)
add_header(s, "06 · PRINCIPLES", "Data-viz choices, made explicit", 7)
add_footer(s)

principles = [
    ("Data-ink ratio",
     "Tufte: maximise the share of pixels that encode data.",
     ["Removed gridlines from 0.06 → 0.04 opacity",
      "Stat bar trimmed from 6 cells to 4 — dropped non-data labels",
      "Header badge and chart frames removed",
      "Legend cut from 28 to 11 entries (top providers + 'Other')"]),
    ("Dual encoding",
     "Color is not the only channel — works for colorblind viewers.",
     ["Provider encoded by both color AND marker shape",
      "Pareto frontier uses dotted line + cyan accent",
      "Tier separators on rankings use line weight + spacing",
      "Stat values sized 22pt vs labels 10pt — pre-attentive hierarchy"]),
    ("Gestalt grouping",
     "Proximity and similarity tell users what belongs together.",
     ["Filter bar visually attached to chart it controls",
      "Cards share radius, padding, border weight",
      "Same provider color across every tab (255+ models)",
      "Whitespace separates tabs more than borders do"]),
    ("Right chart for the task",
     "Cleveland & McGill: pick the channel humans decode best.",
     ["Position for ranking — bars, not pie charts",
      "Length for magnitude — never area for precise comparison",
      "Log scale where the data spans 1000× — price",
      "Small multiples (faceted scatter) instead of overlapping series"]),
]
gap = Inches(0.2)
cw = (SLIDE_W - Inches(1.2) - gap) / 2
ch = Inches(2.25)
py = Inches(2.0)
for i, (title, sub, bullets) in enumerate(principles):
    r = i // 2
    c = i % 2
    cx = Inches(0.6) + (cw + gap) * c
    cy = py + (ch + gap) * r
    add_rect(s, cx, cy, cw, ch, fill=BG_CARD)
    # Number
    add_text(s, cx + Inches(0.3), cy + Inches(0.2),
             Inches(0.5), Inches(0.4), f"0{i+1}",
             size=11, bold=True, color=ACCENT, tracking=140)
    add_text(s, cx + Inches(0.85), cy + Inches(0.2),
             cw - Inches(1.1), Inches(0.4), title,
             size=15, bold=True, color=TEXT_1)
    add_text(s, cx + Inches(0.85), cy + Inches(0.5),
             cw - Inches(1.1), Inches(0.3), sub,
             size=10, color=TEXT_3)
    # Bullets
    by = cy + Inches(0.85)
    for k, b in enumerate(bullets):
        # tiny accent square as bullet
        add_rect(s, cx + Inches(0.35), by + Inches(0.07) + Inches(0.32) * k,
                 Inches(0.08), Inches(0.08), fill=ACCENT)
        add_text(s, cx + Inches(0.55), by + Inches(0.32) * k,
                 cw - Inches(0.8), Inches(0.3), b,
                 size=10, color=TEXT_2)


# ── Slide 8 — LOCAL & HYBRID ───────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
set_bg(s)
add_header(s, "07 · LOCAL + HYBRID", "Open weights and the Claude Code stack", 8)
add_footer(s)

add_text(s, Inches(0.6), Inches(1.85), Inches(12), Inches(0.5),
         "The Run Local tab models open-weight inference. The Agent Stack tab uses that to recommend hybrid setups for tools like Claude Code.",
         size=13, color=TEXT_2)

# Left: VRAM formula card
lx = Inches(0.6); ly = Inches(2.5); lw = Inches(6.2); lh = Inches(4.2)
add_rect(s, lx, ly, lw, lh, fill=BG_CARD)
add_text(s, lx + Inches(0.3), ly + Inches(0.25), lw - Inches(0.6), Inches(0.3),
         "RUN LOCAL  ·  HARDWARE FIT", size=9, bold=True, color=ACCENT, tracking=140)
add_text(s, lx + Inches(0.3), ly + Inches(0.6), lw - Inches(0.6), Inches(0.5),
         "Will it fit on my GPU?", size=20, bold=True, color=TEXT_1)

add_text(s, lx + Inches(0.3), ly + Inches(1.25), lw - Inches(0.6), Inches(0.3),
         "VRAM  =  params × bytes_per_weight × 1.18", size=12, bold=True, color=ACCENT)
add_text(s, lx + Inches(0.3), ly + Inches(1.55), lw - Inches(0.6), Inches(0.3),
         "Speed  =  bandwidth / model_size × efficiency", size=12, bold=True, color=ACCENT)

add_text(s, lx + Inches(0.3), ly + Inches(2.1), lw - Inches(0.6), Inches(0.3),
         "QUANTIZATION LEVELS", size=8, bold=True, color=TEXT_3, tracking=140)
quant = [("FP16","2.0 B/w"),("Q8","1.0"),("Q5","0.625"),("Q4","0.5"),("Q3","0.375"),("Q2","0.25")]
qw = (lw - Inches(0.6)) / 6
for i, (q, b) in enumerate(quant):
    qx = lx + Inches(0.3) + qw * i
    add_text(s, qx, ly + Inches(2.4), qw, Inches(0.3),
             q, size=12, bold=True, color=TEXT_1)
    add_text(s, qx, ly + Inches(2.65), qw, Inches(0.3),
             b, size=9, color=TEXT_3)

add_text(s, lx + Inches(0.3), ly + Inches(3.15), lw - Inches(0.6), Inches(0.3),
         "EFFICIENCY BY BACKEND", size=8, bold=True, color=TEXT_3, tracking=140)
add_text(s, lx + Inches(0.3), ly + Inches(3.42), lw - Inches(0.6), Inches(0.3),
         "Apple MLX 0.82  ·  NVIDIA / Qualcomm 0.55  ·  AMD / Intel 0.50  ·  CPU 0.30",
         size=11, color=TEXT_2)

# Right: Hybrid stack diagram
rx = Inches(7.0); ry = Inches(2.5); rw = Inches(5.75); rh = Inches(4.2)
add_rect(s, rx, ry, rw, rh, fill=BG_CARD)
add_text(s, rx + Inches(0.3), ry + Inches(0.25), rw - Inches(0.6), Inches(0.3),
         "AGENT STACK  ·  CLAUDE CODE / OPENCLAW", size=9, bold=True, color=ACCENT, tracking=140)
add_text(s, rx + Inches(0.3), ry + Inches(0.6), rw - Inches(0.6), Inches(0.5),
         "Three-tier hybrid workflow", size=20, bold=True, color=TEXT_1)

tiers = [
    ("REASONING",  "Claude Opus 4.6",      "API · cloud",   "Plans, decomposes, delegates",  PURPLE),
    ("BALANCED",   "Claude Sonnet 4.6",    "API · cloud",   "Codes, edits, writes",          ACCENT),
    ("FAST",       "Qwen3-30B-A3B (Q4)",   "Local · M-series", "High-volume sub-tasks",     GREEN),
]
ty = ry + Inches(1.3)
th = Inches(0.85)
for i, (tier, model, where, role, col) in enumerate(tiers):
    yy = ty + (th + Inches(0.08)) * i
    add_rect(s, rx + Inches(0.3), yy, rw - Inches(0.6), th, fill=BG_ELEVATED)
    add_rect(s, rx + Inches(0.3), yy, Inches(0.06), th, fill=col)
    add_text(s, rx + Inches(0.5), yy + Inches(0.13), Inches(1.2), Inches(0.3),
             tier, size=8, bold=True, color=col, tracking=140)
    add_text(s, rx + Inches(0.5), yy + Inches(0.38), Inches(3.0), Inches(0.3),
             model, size=12, bold=True, color=TEXT_1)
    add_text(s, rx + Inches(0.5), yy + Inches(0.6), Inches(3.0), Inches(0.3),
             role, size=9, color=TEXT_3)
    add_text(s, rx + rw - Inches(1.5), yy + Inches(0.3), Inches(1.2), Inches(0.3),
             where, size=10, color=TEXT_2, align=PP_ALIGN.RIGHT)


# ── Slide 9 — KEY FINDINGS ─────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
set_bg(s)
add_header(s, "08 · WHAT WE FOUND", "Three patterns the data exposes", 9)
add_footer(s)

findings = [
    ("01",
     "Price ≠ Quality",
     "Pearson r between log10(price) and intelligence is roughly 0.45 — moderate. Half the variance in price is unexplained by capability alone.",
     "Cheap-and-strong models exist on every Pareto step.",
     ACCENT),
    ("02",
     "The frontier moves monthly",
     "Across 30+ daily snapshots, the Pareto frontier has shifted left (cheaper) more often than up (smarter) — price compression dominates.",
     "Recompute every model decision quarterly, minimum.",
     PURPLE),
    ("03",
     "Open weights are catching up",
     "On Apple Silicon at Q4, several open models now sit within 10 points of mid-tier API offerings — for $0 marginal cost.",
     "Hybrid stacks are no longer a hobbyist trick.",
     GREEN),
]
gap = Inches(0.2)
cw = (SLIDE_W - Inches(1.2) - gap * 2) / 3
ch = Inches(4.6)
fy = Inches(2.0)
for i, (num, title, body, takeaway, col) in enumerate(findings):
    cx = Inches(0.6) + (cw + gap) * i
    add_rect(s, cx, fy, cw, ch, fill=BG_CARD)
    add_text(s, cx + Inches(0.3), fy + Inches(0.3), cw - Inches(0.6),
             Inches(0.5), num, size=36, bold=True, color=col)
    add_text(s, cx + Inches(0.3), fy + Inches(1.05), cw - Inches(0.6),
             Inches(0.5), title, size=18, bold=True, color=TEXT_1)
    add_text(s, cx + Inches(0.3), fy + Inches(1.7), cw - Inches(0.6),
             Inches(2.0), body, size=12, color=TEXT_2)
    # Divider
    add_line(s, cx + Inches(0.3), fy + Inches(3.55),
             cx + cw - Inches(0.3), fy + Inches(3.55), color=BORDER)
    add_text(s, cx + Inches(0.3), fy + Inches(3.7), cw - Inches(0.6),
             Inches(0.3), "TAKEAWAY", size=8, bold=True, color=TEXT_3, tracking=140)
    add_text(s, cx + Inches(0.3), fy + Inches(3.95), cw - Inches(0.6),
             Inches(0.5), takeaway, size=11, bold=True, color=col)


# ── Slide 10 — DEMO + REFERENCES ───────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
set_bg(s)

# Big "DEMO" word, dashboard-style
add_chip(s, Inches(0.6), Inches(0.6), "09 · LIVE", color=ACCENT, fill=ACCENT_DIM)
add_text(s, SLIDE_W - Inches(1.2), Inches(0.65), Inches(0.6), Inches(0.3),
         "10 / 10", size=9, color=TEXT_3, align=PP_ALIGN.RIGHT,
         tracking=120, bold=True)

add_text(s, Inches(0.6), Inches(1.4), Inches(12), Inches(2.0),
         "Live Demo", size=80, bold=True, color=TEXT_1)
add_text(s, Inches(0.6), Inches(2.85), Inches(12), Inches(0.6),
         "→ localhost:8050", size=24, color=ACCENT, bold=True)

# Walkthrough cards
items = [
    ("Overview",       "Pareto frontier in motion"),
    ("Rankings",       "Value leaders + tier breaks"),
    ("Compare",        "Radar across 5 models"),
    ("Run Local",      "Pick a GPU, watch fit update"),
    ("Agent Stack",    "Hybrid Claude Code build"),
]
gap = Inches(0.15)
cw = (SLIDE_W - Inches(1.2) - gap * (len(items)-1)) / len(items)
ch = Inches(1.1)
iy = Inches(4.0)
for i, (tab, note) in enumerate(items):
    cx = Inches(0.6) + (cw + gap) * i
    add_rect(s, cx, iy, cw, ch, fill=BG_CARD)
    add_rect(s, cx, iy, cw, Inches(0.05), fill=ACCENT)
    add_text(s, cx + Inches(0.2), iy + Inches(0.2), cw - Inches(0.4),
             Inches(0.3), tab.upper(), size=10, bold=True, color=TEXT_1, tracking=120)
    add_text(s, cx + Inches(0.2), iy + Inches(0.55), cw - Inches(0.4),
             Inches(0.5), note, size=10, color=TEXT_3)

# References strip
ry = Inches(5.6)
add_line(s, Inches(0.6), ry, SLIDE_W - Inches(0.6), ry, color=BORDER)
add_text(s, Inches(0.6), ry + Inches(0.2), Inches(3), Inches(0.3),
         "REFERENCES", size=9, bold=True, color=TEXT_3, tracking=140)
refs = [
    ("Artificial Analysis",      "artificialanalysis.ai  ·  data source for all LLM and image metrics"),
    ("Tufte, E.R.",              "The Visual Display of Quantitative Information  ·  data-ink ratio"),
    ("Cleveland & McGill",       "Graphical perception (1984)  ·  ranking of visual encodings"),
    ("Plotly / Dash docs",       "plotly.com/python  ·  dash.plotly.com"),
]
for i, (name, sub) in enumerate(refs):
    yy = ry + Inches(0.55) + Inches(0.32) * i
    add_text(s, Inches(0.6), yy, Inches(3), Inches(0.3),
             name, size=11, bold=True, color=TEXT_1)
    add_text(s, Inches(3.6), yy, Inches(9.5), Inches(0.3),
             sub, size=11, color=TEXT_3)

# ── Save ────────────────────────────────────────────────────────────────────
out_path = "AI_Frontier_Presentation.pptx"
prs.save(out_path)
print(f"Wrote {out_path}  ·  {len(prs.slides)} slides")
